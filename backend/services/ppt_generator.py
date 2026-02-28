"""
PPT Generation Service using python-pptx
"""
import os
import uuid
import json
from datetime import datetime
from typing import Dict, Any, List, Optional
from pathlib import Path
import asyncio

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from config import settings
from models.schemas import PPTTheme


class PPTGenerator:
    """PPT file generator"""
    
    # Theme color schemes
    THEME_COLORS = {
        "business": {
            "primary": (0, 51, 102),      # Navy blue
            "secondary": (0, 102, 204),   # Light blue
            "accent": (255, 153, 0),      # Orange
            "text": (51, 51, 51),         # Dark gray
            "background": (255, 255, 255) # White
        },
        "education": {
            "primary": (34, 139, 34),     # Forest green
            "secondary": (85, 107, 47),    # Olive
            "accent": (255, 215, 0),      # Gold
            "text": (33, 33, 33),         # Dark gray
            "background": (250, 250, 250) # Off white
        },
        "creative": {
            "primary": (128, 0, 128),     # Purple
            "secondary": (255, 0, 128),   # Pink
            "accent": (0, 255, 255),      # Cyan
            "text": (33, 33, 33),         # Dark gray
            "background": (255, 255, 255)
        },
        "minimal": {
            "primary": (0, 0, 0),         # Black
            "secondary": (128, 128, 128), # Gray
            "accent": (0, 0, 0),          # Black
            "text": (33, 33, 33),         # Dark gray
            "background": (255, 255, 255)
        },
        "tech": {
            "primary": (0, 100, 200),     # Tech blue
            "secondary": (0, 50, 100),    # Dark blue
            "accent": (0, 255, 150),      # Neon green
            "text": (240, 240, 240),      # Light gray
            "background": (20, 20, 40)    # Dark background
        }
    }
    
    def __init__(self):
        self.output_dir = Path(settings.OUTPUT_DIR)
        self.output_dir.mkdir(parents=True, exist_ok=True)
    
    def _rgb_to_pptx(self, rgb: tuple) -> RgbColor:
        """Convert RGB tuple to PPTX color"""
        return RgbColor(rgb[0], rgb[1], rgb[2])
    
    def _apply_title_slide(
        self, 
        slide, 
        title: str, 
        subtitle: str,
        colors: Dict[str, tuple]
    ):
        """Apply title slide design"""
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Style title
        title_frame = title_shape.text_frame
        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(44)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self._rgb_to_pptx(colors["primary"])
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Add subtitle
        if len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_frame = subtitle_shape.text_frame
            for paragraph in subtitle_frame.paragraphs:
                paragraph.font.size = Pt(24)
                paragraph.font.color.rgb = self._rgb_to_pptx(colors["secondary"])
                paragraph.alignment = PP_ALIGN.CENTER
    
    def _apply_content_slide(
        self,
        slide,
        title: str,
        bullets: List[str],
        colors: Dict[str, tuple]
    ):
        """Apply content slide design"""
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Style title
        title_frame = title_shape.text_frame
        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self._rgb_to_pptx(colors["primary"])
        
        # Add content
        content_shape = slide.placeholders[1]
        text_frame = content_shape.text_frame
        text_frame.clear()
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.font.size = Pt(20)
            p.font.color.rgb = self._rgb_to_pptx(colors["text"])
            p.space_before = Pt(12)
    
    def _apply_two_column_slide(
        self,
        slide,
        title: str,
        left_content: str,
        right_content: str,
        colors: Dict[str, tuple]
    ):
        """Apply two-column slide design"""
        title_shape = slide.shapes.title
        title_shape.text = title
        
        title_frame = title_shape.text_frame
        for paragraph in title_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self._rgb_to_pptx(colors["primary"])
        
        # Create two columns
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(4)
        height = Inches(4)
        
        # Left column
        left_box = slide.shapes.add_textbox(left, top, width, height)
        left_frame = left_box.text_frame
        left_frame.text = left_content
        for p in left_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.color.rgb = self._rgb_to_pptx(colors["text"])
        
        # Right column
        right_box = slide.shapes.add_textbox(left + width + Inches(0.5), top, width, height)
        right_frame = right_box.text_frame
        right_frame.text = right_content
        for p in right_frame.paragraphs:
            p.font.size = Pt(18)
            p.font.color.rgb = self._rgb_to_pptx(colors["text"])
    
    async def generate(
        self,
        topic: str,
        outline: List[Dict[str, Any]],
        theme: str = "business",
        language: str = "zh"
    ) -> Dict[str, Any]:
        """Generate PPT file from outline"""
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        colors = self.THEME_COLORS.get(theme, self.THEME_COLORS["business"])
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        self._apply_title_slide(
            slide, 
            topic, 
            f"Created: {datetime.now().strftime('%Y-%m-%d')}",
            colors
        )
        
        # Content slides
        for page in outline:
            page_num = page.get("page_num", 1)
            
            if page_num == 1:
                continue  # Skip title page
            
            # Use different layouts based on content
            if len(page.get("bullets", [])) <= 3:
                # Use two-column layout for detailed content
                content_slide_layout = prs.slide_layouts[5]  # Title only
                slide = prs.slides.add_slide(content_slide_layout)
                
                left = page.get("content", "")
                right = "\n\n".join(page.get("bullets", []))
                self._apply_two_column_slide(
                    slide,
                    page.get("title", f"Page {page_num}"),
                    left,
                    right,
                    colors
                )
            else:
                # Use bullet layout
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                self._apply_content_slide(
                    slide,
                    page.get("title", f"Page {page_num}"),
                    page.get("bullets", []),
                    colors
                )
        
        # Thank you slide
        thank_you_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(thank_you_layout)
        title_shape = slide.shapes.title
        title_shape.text = "谢谢！" if language == "zh" else "Thank You!"
        
        title_frame = title_shape.text_frame
        for p in title_frame.paragraphs:
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = self._rgb_to_pptx(colors["primary"])
            p.alignment = PP_ALIGN.CENTER
        
        # Save file
        task_id = str(uuid.uuid4())
        filename = f"ppt_{task_id}.pptx"
        filepath = self.output_dir / filename
        
        prs.save(str(filepath))
        
        return {
            "task_id": task_id,
            "filename": filename,
            "filepath": str(filepath),
            "file_size": os.path.getsize(filepath)
        }
    
    def get_file_path(self, filename: str) -> Optional[Path]:
        """Get full file path"""
        filepath = self.output_dir / filename
        if filepath.exists():
            return filepath
        return None


# Singleton instance
ppt_generator = PPTGenerator()
